"""
Usage:
  hoge.py -h
  hoge.py (--input-path <input-path>) (--output-path <output-path>) (--column <column>) (--start-x <start-x>) (--start-y <start-y>) (--offset-x <offset-x>) (--offset-y <offset-y>) (--height <height>)

Options:
  -h --help                      Show this screen.
  --input-path <input-path>      Input image folder.
  --output-path <output-path>    Output pptx file path.
  --column <column>              Image column.
  --start-x <start-x>            Image start x.
  --start-y <start-y>            Image start y.
  --offset-x <offset-x>          Image offset x.
  --offset-y <offset-y>          Image offst y.
  --height <height>              Image height.
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from glob import glob
from docopt import docopt


def Centis(length):
	centi = Inches(length / 2.54)
	return centi


if __name__ == '__main__':

	args = docopt(__doc__)

	# Power Point のインスタンス
	ppt = Presentation()

	# 白紙レイアウト
	blank_slide_layout = ppt.slide_layouts[6]

	# 画像の配置を開始する位置
	start_x = Centis(float(args["--start-x"]))
	start_y = Centis(float(args["--start-y"]))

	# 画像の配置間隔
	offset_x = Centis(float(args["--offset-x"]))
	offset_y = Centis(float(args["--offset-y"]))
	image_height = Centis(float(args["--height"]))

	# スライドに貼り付けるすべての画像ファイル
	files = glob(args["--input-path"])
	fileCount = len(files)

	# 作成するスライドの数
	slide_count = int(len(files) / 10) + 1

	column = int(args["--column"])

	for i in range(slide_count):

		# スライドを作成
		slide = ppt.slides.add_slide(blank_slide_layout)

		for j in range(10):
			# 貼り付ける画像の取得
			index = i * 10 + j
			if fileCount <= index:
				continue
			file = files[index]

			# 画像を貼り付ける位置
			x = start_x + j % column * offset_x
			y = start_y + int(j / column) * offset_y

			# 画像を貼り付ける
			slide.shapes.add_picture(file, x, y, height=image_height)

	# .pptx ファイルを保存
	ppt.save(args["--output-path"])
